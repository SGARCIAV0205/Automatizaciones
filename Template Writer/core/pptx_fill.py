from typing import Dict
from pptx import Presentation


def fill_pptx_template(template_path: str, output_path: str, mapping: Dict[str, str]) -> None:
    prs = Presentation(template_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    text = run.text
                    for k, v in mapping.items():
                        token = "{{" + k + "}}"
                        if token in text:
                            text = text.replace(token, v)
                    run.text = text

    prs.save(output_path)

