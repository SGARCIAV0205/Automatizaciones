import re
from typing import List, Set

from docx import Document
from pptx import Presentation

PLACEHOLDER_RE = re.compile(r"\{\{\s*([A-Za-z0-9_]+)\s*\}\}")


def _extract_from_text(text: str) -> Set[str]:
    return set(m.group(1) for m in PLACEHOLDER_RE.finditer(text or ""))


def extract_placeholders_docx(path: str) -> List[str]:
    doc = Document(path)
    found: Set[str] = set()

    for p in doc.paragraphs:
        found |= _extract_from_text(p.text)

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    found |= _extract_from_text(p.text)

    return sorted(found)


def extract_placeholders_pptx(path: str) -> List[str]:
    prs = Presentation(path)
    found: Set[str] = set()

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            for paragraph in tf.paragraphs:
                found |= _extract_from_text(paragraph.text)

    return sorted(found)
