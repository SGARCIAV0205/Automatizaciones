# ==========================================
# src/build_report.py — versión FINAL con placeholders
# ==========================================

import os
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from .utils import load_cfg, ensure_dirs, today_iso
from .brechas_plot import compute_gaps

# ------------------------------------------
# Helpers de formato
# ------------------------------------------
def rgb(hexstr):
    h = hexstr.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def replace_text(slide, key, value, color=None, size=None, bold=False, font="Space Grotesk"):
    """Reemplaza texto dentro de placeholders existentes."""
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame") or not shape.has_text_frame:
            continue
        if key in shape.text:
            shape.text = shape.text.replace(key, value)
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = font
                    if color: r.font.color.rgb = color
                    if size: r.font.size = Pt(size)
                    r.font.bold = bold

def prev_period(p):
    y, m = map(int, p.split("-"))
    m -= 1
    if m == 0:
        y -= 1
        m = 12
    return f"{y}-{m:02d}"

def load_scores(periodo):
    path = f"data/processed/radar_scores_{periodo}.csv"
    return pd.read_csv(path) if os.path.exists(path) else None

# ------------------------------------------
# Narrativa
# ------------------------------------------
def auto_resumen_ejecutivo(scores_cur: pd.DataFrame, axes, news_n: int) -> str:
    leader = scores_cur.iloc[0]["provider"]
    leader_sc = scores_cur.iloc[0]["score_ponderado"]
    spread = leader_sc - scores_cur["score_ponderado"].iloc[-1]

    ub = scores_cur[scores_cur["provider"] == "UBIMIA"]
    if not ub.empty:
        ub_row = ub.iloc[0]
        fortalezas = ", ".join(pd.Series(axes)[np.argsort(-ub_row[axes].values)[:2]])
        debiles = ", ".join(pd.Series(axes)[np.argsort( ub_row[axes].values)[:2]])
    else:
        fortalezas, debiles = "—", "—"

    return (
        f"{leader} lidera el periodo con {leader_sc:.1f} puntos (spread {spread:.1f}). "
        f"Se analizaron {len(scores_cur)} competidores y {news_n} noticias relevantes. "
        f"Para UBIMIA, principales fortalezas: {fortalezas}; focos de mejora: {debiles}."
    )

def analisis_desempeno(scores_cur: pd.DataFrame) -> str:
    leader = scores_cur.iloc[0]["provider"]
    leader_sc = scores_cur.iloc[0]["score_ponderado"]
    median_sc = scores_cur["score_ponderado"].median()
    p75 = scores_cur["score_ponderado"].quantile(0.75)
    p25 = scores_cur["score_ponderado"].quantile(0.25)
    return (
        f"{leader} encabeza el ranking con {leader_sc:.1f} puntos, por encima de la mediana sectorial ({median_sc:.1f}). "
        f"La dispersión de desempeño se concentra entre {p25:.1f} y {p75:.1f} puntos, "
        f"lo que sugiere un entorno competitivo estable. "
        f"Los resultados del periodo confirman posiciones consistentes y brechas claras entre líderes y retadores."
    )

PRIORITY_TEMPLATES = {
    "Integraciones": "Ampliar catálogo de integraciones (API/partners prioritarios) y fortalecer conectores críticos.",
    "Implementación": "Reducir tiempos de despliegue y estandarizar playbooks.",
    "IA/Automatización": "Acelerar features de automatización con métricas claras.",
    "Seguridad/Compliance": "Reforzar certificaciones y comunicación de controles.",
    "Pricing/Valor": "Optimizar empaquetados y condiciones comerciales.",
    "Soporte/SLAs": "Fortalecer tiempos de respuesta y SLA adherence.",
    "Tracción de mercado": "Incrementar visibilidad y adquisición en cuentas objetivo.",
}

def build_top3_priorities(scores_cur, axes, focal="UBIMIA"):
    gaps = compute_gaps(scores_cur, axes, focal=focal)
    gaps = gaps.sort_values("gap", ascending=False)
    ejes_brecha = [e for e in gaps[gaps["gap"] > 0]["eje"].tolist()]

    fallback = [
        "Consolidar propuesta de valor transversal.",
        "Profundizar integraciones de alto impacto.",
        "Comunicar avances técnicos a mercado.",
    ]
    out = []

    for eje in ejes_brecha:
        if len(out) >= 3: break
        txt = PRIORITY_TEMPLATES.get(eje, f"Mejorar desempeño en {eje}.")
        out.append(txt)

    for fb in fallback:
        if len(out) >= 3: break
        out.append(fb)

    return out[:3]

# ------------------------------------------
# Noticias
# ------------------------------------------
def clean_news(df):
    if df.empty: return df
    df = df[~df["fuente"].fillna("").str.contains("google", case=False)]
    df = df[~df["titular"].fillna("").str.contains(r"^\[WARN\]", case=False)]
    df["key"] = df["empresa"].fillna("") + "||" + df["titular"].fillna("")
    return df.drop_duplicates("key")

def chunk_list(items, n):
    for i in range(0, len(items), n):
        yield items[i:i+n]

# ------------------------------------------
# Build
# ------------------------------------------
def main():
    cfg = load_cfg()
    ensure_dirs()

    periodo = cfg["periodo"]
    brand = cfg["brand"]
    QB = rgb(brand["quantum_blue"])
    MS = rgb(brand["mint_signal"])

    axes_cfg = cfg["axes"]
    axes = [a["name"] for a in axes_cfg]

    scores_cur = load_scores(periodo)
    scores_cur = scores_cur.sort_values("score_ponderado", ascending=False).reset_index(drop=True)

    p_m1 = prev_period(periodo)
    scores_prev = load_scores(p_m1)

    radar_png    = f"data/processed/radar_{periodo}.png"
    brechas_png  = f"data/processed/brechas_{periodo}.png"
    histor_png   = f"data/processed/historico_{periodo}.png"

    # ---------------- Noticias ----------------
    news_path = (
        f"data/processed/noticias_enriquecidas_{periodo}.csv"
        if os.path.exists(f"data/processed/noticias_enriquecidas_{periodo}.csv")
        else f"data/processed/noticias_{periodo}.csv"
    )
    news = pd.read_csv(news_path) if os.path.exists(news_path) else pd.DataFrame(columns=["fecha","empresa","titular","fuente"])
    news = clean_news(news)

    news_lines = [
        f"- {r.get('fecha','')} · {r.get('empresa','')}: {r.get('titular','')} ({r.get('fuente','')})"
        for _, r in news.iterrows()
    ]

    # ---------------- Plantilla ----------------
    prs = Presentation("reports/plantilla_radar_UBIMIA.pptx")

    # ------------------------------------------
    # SLIDE 1 — Portada
    # ------------------------------------------
    s0 = prs.slides[0]
    replace_text(s0, "{{RESUMEN_EJECUTIVO}}", auto_resumen_ejecutivo(scores_cur, axes, len(news)), color=RGBColor(255,255,255), size=16)

    # ------------------------------------------
    # SLIDE 2 — Radar + explicaciones
    # ------------------------------------------
    s1 = prs.slides[1]
    replace_text(s1, "{{EXPLICACION_RADAR}}",
                 "Objetivo: comparar capacidades clave entre competidores.\n"
                 "Escala 1–5 por eje con ponderación.\n"
                 "Lectura: figuras grandes y equilibradas indican solidez.\n"
                 "Brechas señalan áreas de oportunidad prioritaria.",
                 color=QB, size=12)

    replace_text(s1, "{{DESCRIPCION_EJES}}",
                 "Integraciones · Implementación · IA/Automatización · Seguridad\n"
                 "Pricing/Valor · Soporte/SLAs · Tracción de mercado",
                 color=QB, size=11)

    # Insertar radar
    if os.path.exists(radar_png):
        s1.shapes.add_picture(radar_png, Inches(0.3), Inches(1.5), width=Inches(6.0))

    # ------------------------------------------
    # SLIDE 3 — Ranking + análisis + deltas
    # ------------------------------------------
    s2 = prs.slides[2]
    ranking_text = "\n".join([f"{i+1}. {r['provider']} — {r['score_ponderado']:.1f}" for i, r in scores_cur.iterrows()])
    replace_text(s2, "{{TABLA_RANKING}}", ranking_text, color=QB, size=16)

    replace_text(s2, "{{TEXTO_ANALISIS_DESEMPENO}}", analisis_desempeno(scores_cur), color=QB, size=14)

    # Deltas
    deltas_txt = "Sin histórico del mes anterior."
    if scores_prev is not None:
        prev = scores_prev.rename(columns={"score_ponderado":"score_prev"})
        cur  = scores_cur.rename(columns={"score_ponderado":"score_cur"})
        merged = cur.merge(prev[["provider","score_prev"]], on="provider", how="left")
        merged["Δ"] = (merged["score_cur"] - merged["score_prev"]).round(1)
        deltas_txt = "\n".join([f"{r.provider}: {r._asdict()['Δ']:+.1f}" for r in merged.itertuples(index=False)])
    replace_text(s2, "{{TABLA_DELTAS}}", deltas_txt, color=QB, size=12)

    # ------------------------------------------
    # SLIDE 4 — Brechas + prioridades
    # ------------------------------------------
    s3 = prs.slides[3]

    if os.path.exists(brechas_png):
        s3.shapes.add_picture(brechas_png, Inches(0.3), Inches(1.5), width=Inches(7.5))

    top3 = build_top3_priorities(scores_cur, axes)
    replace_text(s3, "{{TOP3_PRIORIDADES}}",
                 "Top 3 prioridades\n" + "\n".join([f"• {x}" for x in top3]),
                 color=QB, size=14)

    # ------------------------------------------
    # SLIDE 5 — Histórico
    # ------------------------------------------
    s4 = prs.slides[4]
    if os.path.exists(histor_png):
        s4.shapes.add_picture(histor_png, Inches(0.3), Inches(1.5), width=Inches(11.5))
    else:
        replace_text(s4, "{{TEXTO_HISTORICO}}", "Aún no hay histórico suficiente.", color=QB, size=14)

    # ------------------------------------------
    # SLIDE 6 — Noticias
    # ------------------------------------------
    s5 = prs.slides[5]
    first_page_news = "\n".join(news_lines[:12]) if news_lines else "Sin noticias relevantes en el periodo."
    replace_text(s5, "{{TABLA_NOTICIAS}}", first_page_news, color=QB, size=14)

    # Paginar noticias extra
    chunks = list(chunk_list(news_lines, 12))
    if len(chunks) > 1:
        blank = prs.slide_layouts[6]
        for i, chunk in enumerate(chunks[1:], start=2):
            slide = prs.slides.add_slide(blank)
            textbox = slide.shapes.add_textbox(Inches(0.4), Inches(0.5), Inches(12.0), Inches(6.0))
            textbox.text = "\n".join(chunk)
            for p in textbox.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = "Space Grotesk"
                    r.font.color.rgb = QB
                    r.font.size = Pt(14)

    # ------------------------------------------
    # Guardar reporte
    # ------------------------------------------
    out_pptx = f"reports/out/reporte_radar_UBIMIA_auto_v3_{periodo}.pptx"
    prs.save(out_pptx)
    print(f"[OK] Reporte PPTX generado: {out_pptx}")

if __name__ == "__main__":
    main()
