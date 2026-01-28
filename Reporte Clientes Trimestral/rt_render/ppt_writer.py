# render/ppt_writer.py

from copy import deepcopy
from pathlib import Path
from pptx import Presentation

from rt_utils.placeholders import CLIENT_PLACEHOLDERS


# ============================================================
# Helpers
# ============================================================

def _replace_text_in_shape(shape, mapping: dict):
    if not shape.has_text_frame:
        return

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            for key, value in mapping.items():
                placeholder = "{" + key + "}"
                if placeholder in run.text:
                    run.text = run.text.replace(
                        placeholder,
                        value if value is not None else ""
                    )


def _replace_text_in_slide(slide, mapping: dict):
    for shape in slide.shapes:
        _replace_text_in_shape(shape, mapping)


def _duplicate_slide(prs, slide):
    slide_layout = slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    for shape in slide.shapes:
        new_shape = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(
            new_shape, 'p:extLst'
        )

    return new_slide


def _delete_slide_by_index(prs, index: int):
    """
    Elimina un slide de forma segura usando su índice.
    """
    slide_id_list = prs.slides._sldIdLst
    slide_id_list.remove(slide_id_list[index])


# ============================================================
# API pública
# ============================================================

def generate_ppt(
    global_payload: dict,
    clients_payload: list,
    template_path: Path,
    output_path: Path
):
    prs = Presentation(str(template_path))

    # --------------------------------------------------------
    # 1. Reemplazos globales
    # --------------------------------------------------------

    global_mapping = {
        "PERIODO": global_payload.get("PERIODO", ""),
        "LISTA_CLIENTES": global_payload.get("LISTA_CLIENTES", ""),
        "REPORT_TITLE": global_payload.get("REPORT_TITLE", "")
    }

    for slide in prs.slides:
        _replace_text_in_slide(slide, global_mapping)

    # --------------------------------------------------------
    # 2. Identificar slides base de cliente (por índice)
    # --------------------------------------------------------

    client_slide_indices = []

    for i, slide in enumerate(prs.slides):
        slide_text = " ".join(
            shape.text for shape in slide.shapes if shape.has_text_frame
        )
        if "{CLIENTE}" in slide_text:
            client_slide_indices.append(i)

    # Guardamos una copia de los slides base
    base_client_slides = [prs.slides[i] for i in client_slide_indices]

    # --------------------------------------------------------
    # 3. Generar slides por cliente
    # --------------------------------------------------------

    for client_payload in clients_payload:
        mapping = {
            key: client_payload.get(key, "")
            for key in CLIENT_PLACEHOLDERS.keys()
        }

        for base_slide in base_client_slides:
            new_slide = _duplicate_slide(prs, base_slide)
            _replace_text_in_slide(new_slide, mapping)

    # --------------------------------------------------------
    # 4. Eliminar slides base (en orden inverso)
    # --------------------------------------------------------

    for idx in sorted(client_slide_indices, reverse=True):
        _delete_slide_by_index(prs, idx)

    # --------------------------------------------------------
    # 5. Guardar archivo final
    # --------------------------------------------------------

    prs.save(str(output_path))
