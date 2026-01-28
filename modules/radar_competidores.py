from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

def build_competitor_radar(
    competitors: list,
    news_window_days: int,
    theme: str,
    title: str,
    output_path: Path,
    output_format: str = "pptx"
) -> bool:
    """
    Wrapper base de ejemplo (usa python-pptx).
    Conecta tu scraper/generador real desde la página con el cargador dinámico.
    """
    output_path.parent.mkdir(exist_ok=True, parents=True)

    prs = Presentation()

    # Portada
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tx = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(1.5))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    sub = tf.add_paragraph()
    sub.text = f"Tema: {theme} | Ventana: {news_window_days} días"
    sub.level = 1

    # Contenido por competidor (mock)
    for comp in competitors:
        s = prs.slides.add_slide(prs.slide_layouts[5])
        tbox = s.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(0.8))
        tf2 = tbox.text_frame
        tf2.text = comp
        tf2.paragraphs[0].font.size = Pt(32)
        tf2.paragraphs[0].font.bold = True

        body = s.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        bf = body.text_frame
        bf.text = "Principales titulares (ejemplo - integra tu scraper aquí):"
        for i in range(1, 4):
            pp = bf.add_paragraph()
            pp.text = f"- Noticia {i} de {comp} (últimos {news_window_days} días)"
            pp.level = 1
            pp.font.size = Pt(16)

    prs.save(output_path)
    return True
